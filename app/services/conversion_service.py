from __future__ import annotations

import re
from pathlib import Path

from app.services.pdf_service import PDFService


class ConversionService:
    @staticmethod
    def _write_lines_to_pdf(lines: list[str], output_path: str, title: str = "Document Export") -> str:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        pdf = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        y = height - 48
        pdf.setFont("Helvetica-Bold", 14)
        pdf.drawString(40, y, title)
        y -= 26
        pdf.setFont("Helvetica", 10)
        for line in lines:
            text = (line or "").strip()
            if not text:
                y -= 8
                if y < 40:
                    pdf.showPage()
                    y = height - 40
                    pdf.setFont("Helvetica", 10)
                continue
            chunks = [text[index : index + 100] for index in range(0, len(text), 100)]
            for chunk in chunks:
                if y < 40:
                    pdf.showPage()
                    y = height - 40
                    pdf.setFont("Helvetica", 10)
                pdf.drawString(40, y, chunk)
                y -= 14
        pdf.save()
        return output_path

    @staticmethod
    def pdf_to_word(input_path: str, output_path: str) -> str:
        from docx import Document

        text_file = Path(output_path).with_suffix(".txt")
        PDFService.pdf_to_text(input_path, str(text_file))
        document = Document()
        document.add_heading("PDF Export", level=1)
        content = text_file.read_text(encoding="utf-8")
        for paragraph in content.split("\n\n"):
            if paragraph.strip():
                document.add_paragraph(paragraph.strip())
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        document.save(output_path)
        return output_path

    @staticmethod
    def pdf_to_powerpoint(input_path: str, output_path: str, working_dir: str) -> str:
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        image_paths = PDFService.pdf_to_images(input_path, working_dir, image_format="png")
        for image_path in image_paths:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(image_path, 0, 0, width=presentation.slide_width, height=presentation.slide_height)
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        return output_path

    @staticmethod
    def pdf_to_excel(input_path: str, output_path: str) -> str:
        import pdfplumber
        from openpyxl import Workbook

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Extracted PDF"
        row_index = 1
        with pdfplumber.open(input_path) as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()
                sheet.cell(row=row_index, column=1, value=f"Page {page_number}")
                row_index += 1
                if tables:
                    for table in tables:
                        for row in table:
                            for column_index, cell in enumerate(row or [], start=1):
                                sheet.cell(row=row_index, column=column_index, value=cell)
                            row_index += 1
                else:
                    text = (page.extract_text() or "").splitlines()
                    for line in text:
                        sheet.cell(row=row_index, column=1, value=line)
                        row_index += 1
                row_index += 1
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        workbook.save(output_path)
        return output_path

    @staticmethod
    def pdf_to_html(input_path: str, output_path: str) -> str:
        text_path = Path(output_path).with_suffix(".txt")
        PDFService.pdf_to_text(input_path, str(text_path))
        lines = text_path.read_text(encoding="utf-8", errors="ignore").splitlines()
        escaped_lines = [line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;") for line in lines]
        html = (
            "<!doctype html>\n"
            "<html lang=\"en\"><head><meta charset=\"utf-8\"><title>PDF to HTML</title></head>"
            "<body><pre>\n"
            f"{chr(10).join(escaped_lines)}\n"
            "</pre></body></html>\n"
        )
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        Path(output_path).write_text(html, encoding="utf-8")
        return output_path

    @staticmethod
    def pdf_to_rtf(input_path: str, output_path: str) -> str:
        text_path = Path(output_path).with_suffix(".txt")
        PDFService.pdf_to_text(input_path, str(text_path))
        lines = text_path.read_text(encoding="utf-8", errors="ignore").splitlines()
        escaped = [line.replace("\\", "\\\\").replace("{", "\\{").replace("}", "\\}") for line in lines]
        rtf = "{\\rtf1\\ansi\\deff0\n" + "\\line\n".join(escaped) + "\n}\n"
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        Path(output_path).write_text(rtf, encoding="utf-8")
        return output_path

    @staticmethod
    def resize_image(input_path: str, output_path: str, width: int, height: int) -> str:
        from PIL import Image

        image = Image.open(input_path)
        resized = image.resize((width, height))
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        resized.save(output_path)
        return output_path

    @staticmethod
    def compress_image(input_path: str, output_path: str, quality: int = 75) -> str:
        from PIL import Image

        image = Image.open(input_path).convert("RGB")
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        image.save(output_path, optimize=True, quality=max(20, min(quality, 95)))
        return output_path

    @staticmethod
    def convert_image_format(input_path: str, output_path: str) -> str:
        from PIL import Image

        image = Image.open(input_path)
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        image.save(output_path)
        return output_path

    @staticmethod
    def word_to_pdf(input_path: str, output_path: str) -> str:
        suffix = Path(input_path).suffix.lower()
        lines: list[str] = []
        if suffix == ".docx":
            from docx import Document

            document = Document(input_path)
            lines = [paragraph.text for paragraph in document.paragraphs]
        else:
            lines = Path(input_path).read_text(encoding="utf-8", errors="ignore").splitlines()
        return ConversionService._write_lines_to_pdf(lines, output_path, title="Word to PDF")

    @staticmethod
    def powerpoint_to_pdf(input_path: str, output_path: str) -> str:
        from pptx import Presentation

        lines: list[str] = []
        try:
            presentation = Presentation(input_path)
            for slide_index, slide in enumerate(presentation.slides, start=1):
                lines.append(f"Slide {slide_index}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        for line in shape.text.splitlines():
                            if line.strip():
                                lines.append(f"  {line.strip()}")
                lines.append("")
        except Exception:
            raw_text = Path(input_path).read_text(encoding="utf-8", errors="ignore")
            lines = raw_text.splitlines() or ["Unable to parse binary PPT layout; extracted raw text fallback."]
        return ConversionService._write_lines_to_pdf(lines, output_path, title="PowerPoint to PDF")

    @staticmethod
    def excel_to_pdf(input_path: str, output_path: str) -> str:
        import csv

        from openpyxl import load_workbook

        lines: list[str] = []
        suffix = Path(input_path).suffix.lower()
        if suffix == ".csv":
            with Path(input_path).open("r", encoding="utf-8", errors="ignore", newline="") as handle:
                reader = csv.reader(handle)
                for row in reader:
                    values = [str(cell).strip() for cell in row if cell not in (None, "")]
                    if values:
                        lines.append(" | ".join(values))
        else:
            try:
                workbook = load_workbook(input_path, data_only=True)
            except Exception:
                raw_text = Path(input_path).read_text(encoding="utf-8", errors="ignore")
                lines = raw_text.splitlines() or ["Unable to parse spreadsheet binary; extracted raw text fallback."]
            else:
                for sheet in workbook.worksheets:
                    lines.append(f"[Sheet] {sheet.title}")
                    for row in sheet.iter_rows(min_row=1, max_row=200, values_only=True):
                        values = [str(cell).strip() for cell in row if cell not in (None, "")]
                        if values:
                            lines.append(" | ".join(values))
                    lines.append("")
        return ConversionService._write_lines_to_pdf(lines, output_path, title="Excel to PDF")

    @staticmethod
    def html_to_pdf(input_path: str, output_path: str) -> str:
        html_text = Path(input_path).read_text(encoding="utf-8", errors="ignore")
        body = re.sub(r"<[^>]+>", " ", html_text)
        body = re.sub(r"\s+", " ", body).strip()
        lines = [body[index : index + 120] for index in range(0, len(body), 120)] or ["(empty html)"]
        return ConversionService._write_lines_to_pdf(lines, output_path, title="HTML to PDF")

    @staticmethod
    def text_to_pdf(input_path: str, output_path: str, title: str = "Text to PDF") -> str:
        lines = Path(input_path).read_text(encoding="utf-8", errors="ignore").splitlines()
        return ConversionService._write_lines_to_pdf(lines, output_path, title=title)
